from pptx import Presentation

def parse_pptx(path: str) -> str:
    """Trích xuất toàn bộ text thô từ file PowerPoint (.pptx).
    
    Args:
        path (str): Đường dẫn đến file PPTX.
    
    Returns:
        str: Văn bản thô được trích xuất từ tất cả các slide của PPTX, cách nhau bởi ký tự xuống dòng.
    
    Raises:
        FileNotFoundError: Nếu file không tồn tại tại đường dẫn được cung cấp.
        Exception: Nếu có lỗi xảy ra trong quá trình xử lý file PPTX (ví dụ: file không hợp lệ).
    """
    try:
        # Mở file PPTX bằng python-pptx
        prs = Presentation(path)
        texts = []
        # Duyệt qua từng slide trong file
        for slide in prs.slides:
            # Duyệt qua từng shape trong slide
            for shape in slide.shapes:
                # Kiểm tra nếu shape có văn bản và không rỗng
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
        # Kết hợp tất cả văn bản từ các shape thành một chuỗi
        return "\n".join(texts)
    except FileNotFoundError:
        raise FileNotFoundError(f"File PPTX không tồn tại: {path}")
    except Exception as e:
        raise Exception(f"Lỗi khi xử lý PPTX: {str(e)}")